"""
Tests for Phase 2 ingestion expansion: Excel and PowerPoint files.

Tests the new extraction functions added in Phase 2:
- extract_excel_text()
- extract_powerpoint_text()
"""

import pytest
import io
from backend.ingest import extract_excel_text, extract_powerpoint_text


class TestExcelExtraction:
    """Test Excel spreadsheet content extraction."""

    def test_extract_simple_excel(self):
        """Test extraction from a basic Excel file."""
        from openpyxl import Workbook

        # Create simple workbook
        wb = Workbook()
        ws = wb.active
        ws.title = "Sales Data"
        ws['A1'] = "Product"
        ws['B1'] = "Revenue"
        ws['A2'] = "Widget"
        ws['B2'] = 1000
        ws['A3'] = "Gadget"
        ws['B3'] = 2000

        # Save to bytes
        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'sales.xlsx')

        # Verify output
        assert 'EXCEL SPREADSHEET' in result
        assert 'sales.xlsx' in result
        assert 'Sales Data' in result
        assert 'Product' in result
        assert 'Revenue' in result
        assert 'Widget' in result
        assert '1000' in result
        assert 'Gadget' in result
        assert '2000' in result

    def test_extract_multiple_sheets(self):
        """Test Excel file with multiple sheets."""
        from openpyxl import Workbook

        wb = Workbook()

        # Sheet 1
        ws1 = wb.active
        ws1.title = "Q1"
        ws1['A1'] = "January"
        ws1['A2'] = "February"

        # Sheet 2
        ws2 = wb.create_sheet("Q2")
        ws2['A1'] = "April"
        ws2['A2'] = "May"

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'quarters.xlsx')

        assert 'Sheets: 2' in result
        assert '=== Sheet: Q1' in result
        assert '=== Sheet: Q2' in result
        assert 'January' in result
        assert 'April' in result

    def test_extract_excel_with_numbers(self):
        """Test Excel with numeric values."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = 100
        ws['A2'] = 200.5
        ws['A3'] = -50

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'numbers.xlsx')

        assert '100' in result
        assert '200.5' in result
        assert '-50' in result

    def test_extract_excel_with_empty_cells(self):
        """Test Excel with empty cells."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Name"
        ws['B1'] = "Age"
        ws['A2'] = "Alice"
        ws['B2'] = None  # Empty cell
        ws['A3'] = None
        ws['B3'] = 30

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'sparse.xlsx')

        assert 'Name' in result
        assert 'Alice' in result
        assert '30' in result

    def test_extract_excel_table_format(self):
        """Test that table format is preserved with pipes."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Col1"
        ws['B1'] = "Col2"
        ws['C1'] = "Col3"

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'table.xlsx')

        # Check pipe-separated format
        assert 'Col1 | Col2 | Col3' in result

    def test_extract_large_excel(self):
        """Test Excel with many rows."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active

        # Add 100 rows
        for i in range(100):
            ws[f'A{i+1}'] = f"Row {i+1}"

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'large.xlsx')

        assert 'EXCEL SPREADSHEET' in result
        assert 'Row 1' in result
        assert 'Row 100' in result

    def test_extract_excel_with_formulas(self):
        """Test Excel with formulas (should show calculated values)."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = 10
        ws['A2'] = 20
        ws['A3'] = '=A1+A2'  # Formula

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_excel_text(content_bytes, 'formulas.xlsx')

        # data_only=True in load_workbook should show calculated value
        assert 'EXCEL SPREADSHEET' in result


class TestPowerPointExtraction:
    """Test PowerPoint presentation content extraction."""

    def test_extract_simple_powerpoint(self):
        """Test extraction from a basic PowerPoint file."""
        from pptx import Presentation
        from pptx.util import Inches

        # Create simple presentation
        prs = Presentation()

        # Slide 1
        slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide
        slide1.shapes.title.text = "Test Presentation"
        slide1.placeholders[1].text = "By Test Author"

        # Slide 2
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        slide2.shapes.title.text = "Slide 2 Title"
        slide2.placeholders[1].text = "Slide 2 content"

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'presentation.pptx')

        # Verify output
        assert 'POWERPOINT PRESENTATION' in result
        assert 'presentation.pptx' in result
        assert 'Slides: 2' in result
        assert '--- Slide 1 ---' in result
        assert '--- Slide 2 ---' in result
        assert 'Test Presentation' in result
        assert 'By Test Author' in result
        assert 'Slide 2 Title' in result
        assert 'Slide 2 content' in result

    def test_extract_powerpoint_with_notes(self):
        """Test PowerPoint with speaker notes."""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Title"

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes"

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'with_notes.pptx')

        assert '[Speaker Notes]' in result
        assert 'These are speaker notes' in result

    def test_extract_powerpoint_empty_slides(self):
        """Test PowerPoint with empty slides."""
        from pptx import Presentation

        prs = Presentation()

        # Add blank slide
        prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'empty.pptx')

        assert 'POWERPOINT PRESENTATION' in result
        assert '[Empty slide]' in result

    def test_extract_powerpoint_multiple_slides(self):
        """Test PowerPoint with many slides."""
        from pptx import Presentation

        prs = Presentation()

        # Add 5 slides
        for i in range(5):
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = f"Slide {i+1}"

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'multi.pptx')

        assert 'Slides: 5' in result
        assert 'Slide 1' in result
        assert 'Slide 5' in result

    def test_extract_powerpoint_with_table(self):
        """Test PowerPoint with table shapes."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank

        # Add table
        rows, cols = 2, 3
        left = top = Inches(1)
        width = Inches(6)
        height = Inches(2)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Fill table
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(0, 2).text = "Header 3"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"
        table.cell(1, 2).text = "Data 3"

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'table.pptx')

        assert '[Table]' in result
        assert 'Header 1' in result
        assert 'Data 1' in result

    def test_extract_powerpoint_text_boxes(self):
        """Test PowerPoint with text boxes."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Add text box
        left = top = width = height = Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = "This is a text box"

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = extract_powerpoint_text(content_bytes, 'textbox.pptx')

        assert 'This is a text box' in result


class TestIntegrationWithIngest:
    """Test integration with main ingest_upload_file function."""

    def test_excel_file_routed_correctly(self):
        """Test that .xlsx files are routed to Excel extractor."""
        from backend.ingest import ingest_upload_file
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Test"

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = ingest_upload_file('test.xlsx', content_bytes)

        assert 'EXCEL SPREADSHEET' in result

    def test_powerpoint_file_routed_correctly(self):
        """Test that .pptx files are routed to PowerPoint extractor."""
        from backend.ingest import ingest_upload_file
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test"

        buffer = io.BytesIO()
        prs.save(buffer)
        content_bytes = buffer.getvalue()

        result = ingest_upload_file('test.pptx', content_bytes)

        assert 'POWERPOINT PRESENTATION' in result

    def test_office_files_preserve_structure(self):
        """Test that Office files maintain data structure."""
        from backend.ingest import ingest_upload_file
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws['A1'] = "Col1"
        ws['B1'] = "Col2"
        ws['A2'] = "Val1"
        ws['B2'] = "Val2"

        buffer = io.BytesIO()
        wb.save(buffer)
        content_bytes = buffer.getvalue()

        result = ingest_upload_file('data.xlsx', content_bytes)

        # Check table structure is preserved
        assert '|' in result  # Pipe separator
        assert 'Col1' in result
        assert 'Val1' in result

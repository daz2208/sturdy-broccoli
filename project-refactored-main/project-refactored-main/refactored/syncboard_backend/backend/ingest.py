"""
Real multimodal content ingestion with YouTube transcription.

This module provides actual AI-powered content processing:
- YouTube video transcription using OpenAI Whisper (with audio compression)
- TikTok video processing
- PDF text extraction
- Audio file transcription
- Web article extraction
- Jupyter notebook extraction (.ipynb) - Phase 1
- Source code file processing (Python, JavaScript, etc.) - Phase 1
- ✨ Excel spreadsheet extraction (.xlsx, .xls) - Phase 2
- ✨ PowerPoint presentation extraction (.pptx) - Phase 2

✅ FIXED: Added audio compression to handle files over Whisper's 25MB limit
✨ Phase 1: 40+ programming languages and Jupyter notebooks
✨ Phase 2: Office Suite (Excel, PowerPoint)

Dependencies:
    pip install yt-dlp openai anthropic pypdf beautifulsoup4 requests python-docx
"""

import os
import tempfile
import logging
import subprocess
from typing import Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# Check for required API keys
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    logger.warning("OPENAI_API_KEY not set - YouTube transcription will fail")

# Whisper API file size limit
WHISPER_MAX_SIZE_MB = 25
WHISPER_MAX_SIZE_BYTES = WHISPER_MAX_SIZE_MB * 1024 * 1024


def download_url(url: str) -> str:
    """
    Download and process content from a URL.
    
    Supports:
    - YouTube videos (transcription via Whisper)
    - TikTok videos (transcription via Whisper)
    - Web articles (text extraction)
    - Direct media files
    
    Args:
        url: The URL to process
        
    Returns:
        Extracted text content
        
    Raises:
        Exception: If processing fails
    """
    url_lower = url.lower()
    
    # YouTube video
    if 'youtube.com' in url_lower or 'youtu.be' in url_lower:
        return transcribe_youtube(url)
    
    # TikTok video
    elif 'tiktok.com' in url_lower:
        return transcribe_tiktok(url)
    
    # Regular web page
    else:
        return extract_web_article(url)


def compress_audio_for_whisper(input_path: Path, output_path: Path) -> None:
    """
    Compress audio file to meet Whisper's 25MB limit.
    
    Compression settings optimized for speech transcription:
    - 16kHz sample rate (Whisper's recommended format)
    - Mono channel (speech doesn't need stereo)
    - 64kbps bitrate (sufficient for clear speech)
    
    This typically reduces file size by 50-70% with no quality loss for transcription.
    
    Args:
        input_path: Path to original audio file
        output_path: Path to save compressed audio
        
    Raises:
        Exception: If FFmpeg compression fails
    """
    try:
        subprocess.run([
            'ffmpeg',
            '-i', str(input_path),
            '-ar', '16000',      # 16kHz sample rate (Whisper optimal)
            '-ac', '1',          # Mono audio (sufficient for speech)
            '-b:a', '64k',       # 64kbps bitrate (good quality speech)
            '-y',                # Overwrite output file
            str(output_path)
        ], check=True, capture_output=True, text=True)
        
        original_size = input_path.stat().st_size / (1024 * 1024)  # MB
        compressed_size = output_path.stat().st_size / (1024 * 1024)  # MB
        
        logger.info(
            f"Compressed audio: {original_size:.2f}MB → {compressed_size:.2f}MB "
            f"({100 * (1 - compressed_size/original_size):.1f}% reduction)"
        )
        
    except subprocess.CalledProcessError as e:
        raise Exception(f"Audio compression failed: {e.stderr}")


def chunk_audio_file(audio_path: Path, chunk_duration_seconds: int = 600) -> list[Path]:
    """
    Split audio file into chunks for very long videos.
    
    Only used as fallback if compression still exceeds 25MB limit.
    This typically only happens for videos over 90-120 minutes.
    
    Args:
        audio_path: Path to audio file
        chunk_duration_seconds: Length of each chunk (default 10 minutes)
        
    Returns:
        List of paths to audio chunks
    """
    chunk_pattern = audio_path.parent / f"{audio_path.stem}_chunk_%03d{audio_path.suffix}"
    
    try:
        subprocess.run([
            'ffmpeg',
            '-i', str(audio_path),
            '-f', 'segment',
            '-segment_time', str(chunk_duration_seconds),
            '-c', 'copy',
            '-y',
            str(chunk_pattern)
        ], check=True, capture_output=True, text=True)
        
        # Find all created chunks
        chunks = sorted(audio_path.parent.glob(f"{audio_path.stem}_chunk_*.{audio_path.suffix}"))
        logger.info(f"Split audio into {len(chunks)} chunks")
        return chunks
        
    except subprocess.CalledProcessError as e:
        raise Exception(f"Audio chunking failed: {e.stderr}")


def transcribe_youtube(url: str) -> str:
    """
    Transcribe a YouTube video using OpenAI Whisper.
    
    Process:
    1. Download audio using yt-dlp
    2. Compress audio to meet Whisper's 25MB limit
    3. If still too large, split into chunks
    4. Transcribe with Whisper API
    5. Return transcript with metadata
    
    ✅ FIXED: Now handles videos over 25MB by compressing audio first
    
    Args:
        url: YouTube URL
        
    Returns:
        Full transcript with title and metadata
    """
    try:
        import yt_dlp
        from openai import OpenAI
    except ImportError:
        raise Exception(
            "Missing dependencies. Install with: "
            "pip install yt-dlp openai"
        )
    
    if not OPENAI_API_KEY:
        raise Exception("OPENAI_API_KEY environment variable not set")
    
    logger.info(f"Transcribing YouTube video: {url}")
    
    # Create temporary directory for audio
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        audio_path = temp_path / "audio.mp3"
        compressed_path = temp_path / "audio_compressed.mp3"
        
        # Download audio with yt-dlp
        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': str(audio_path.with_suffix('')),
            'quiet': True,
            'no_warnings': True,
        }
        
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                title = info.get('title', 'Unknown')
                duration = info.get('duration', 0)
                channel = info.get('channel', 'Unknown')
                
            original_size = audio_path.stat().st_size
            logger.info(
                f"Downloaded: {title} ({duration}s, {original_size/(1024*1024):.2f}MB)"
            )
            
        except Exception as e:
            raise Exception(f"Failed to download YouTube video: {e}")
        
        # Check if compression is needed
        if original_size > WHISPER_MAX_SIZE_BYTES:
            logger.info(
                f"Audio file ({original_size/(1024*1024):.2f}MB) exceeds Whisper limit "
                f"({WHISPER_MAX_SIZE_MB}MB). Compressing..."
            )
            compress_audio_for_whisper(audio_path, compressed_path)
            transcription_path = compressed_path
        else:
            logger.info("Audio file within Whisper limit, no compression needed")
            transcription_path = audio_path
        
        # Check if chunking is needed (for very long videos)
        final_size = transcription_path.stat().st_size
        if final_size > WHISPER_MAX_SIZE_BYTES:
            logger.warning(
                f"Even after compression ({final_size/(1024*1024):.2f}MB), "
                f"file exceeds limit. Splitting into chunks..."
            )
            chunks = chunk_audio_file(transcription_path)
            return transcribe_audio_chunks(chunks, title, channel, duration, url)
        
        # Transcribe with Whisper (single file)
        try:
            client = OpenAI(api_key=OPENAI_API_KEY)
            
            logger.info(f"Sending to Whisper API ({final_size/(1024*1024):.2f}MB)...")
            with open(transcription_path, 'rb') as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    response_format="text"
                )
            
            # Format result with metadata
            result = f"""YOUTUBE VIDEO TRANSCRIPT
Title: {title}
Channel: {channel}
Duration: {duration} seconds
URL: {url}

TRANSCRIPT:
{transcript}
"""
            
            logger.info(f"Successfully transcribed {len(transcript)} characters")
            return result
            
        except Exception as e:
            raise Exception(f"Whisper transcription failed: {e}")


def transcribe_audio_chunks(chunks: list[Path], title: str, channel: str, 
                            duration: int, url: str) -> str:
    """
    Transcribe multiple audio chunks and combine results.
    
    Only used for very long videos (90+ minutes) where even compression
    doesn't bring the file under 25MB.
    
    Args:
        chunks: List of audio chunk paths
        title: Video title
        channel: Channel name
        duration: Video duration in seconds
        url: Original video URL
        
    Returns:
        Combined transcript with metadata
    """
    try:
        from openai import OpenAI
        client = OpenAI(api_key=OPENAI_API_KEY)
        
        transcripts = []
        for i, chunk_path in enumerate(chunks, 1):
            logger.info(f"Transcribing chunk {i}/{len(chunks)}...")
            
            with open(chunk_path, 'rb') as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    response_format="text"
                )
            
            transcripts.append(f"[Part {i}]\n{transcript}")
        
        combined_transcript = "\n\n".join(transcripts)
        
        result = f"""YOUTUBE VIDEO TRANSCRIPT
Title: {title}
Channel: {channel}
Duration: {duration} seconds
URL: {url}
Note: Video was split into {len(chunks)} parts for transcription

TRANSCRIPT:
{combined_transcript}
"""
        
        logger.info(f"Successfully transcribed {len(combined_transcript)} characters from {len(chunks)} chunks")
        return result
        
    except Exception as e:
        raise Exception(f"Chunked transcription failed: {e}")


def transcribe_tiktok(url: str) -> str:
    """
    Transcribe a TikTok video.
    
    Similar to YouTube but handles TikTok-specific URL format.
    ✅ Includes audio compression for files over 25MB.
    
    Args:
        url: TikTok URL
        
    Returns:
        Transcript with metadata
    """
    try:
        import yt_dlp
        from openai import OpenAI
    except ImportError:
        raise Exception(
            "Missing dependencies. Install with: "
            "pip install yt-dlp openai"
        )
    
    if not OPENAI_API_KEY:
        raise Exception("OPENAI_API_KEY environment variable not set")
    
    logger.info(f"Transcribing TikTok video: {url}")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        audio_path = temp_path / "audio.mp3"
        compressed_path = temp_path / "audio_compressed.mp3"
        
        ydl_opts = {
            'format': 'bestaudio/best',
            'postprocessors': [{
                'key': 'FFmpegExtractAudio',
                'preferredcodec': 'mp3',
                'preferredquality': '192',
            }],
            'outtmpl': str(audio_path.with_suffix('')),
            'quiet': True,
        }
        
        try:
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                title = info.get('title', 'TikTok Video')
                creator = info.get('creator', 'Unknown')
                
            original_size = audio_path.stat().st_size
            logger.info(f"Downloaded TikTok: {title} ({original_size/(1024*1024):.2f}MB)")
                
        except Exception as e:
            raise Exception(f"Failed to download TikTok video: {e}")
        
        # Compress if needed
        if original_size > WHISPER_MAX_SIZE_BYTES:
            logger.info("Compressing TikTok audio...")
            compress_audio_for_whisper(audio_path, compressed_path)
            transcription_path = compressed_path
        else:
            transcription_path = audio_path
        
        # Transcribe
        try:
            client = OpenAI(api_key=OPENAI_API_KEY)
            
            with open(transcription_path, 'rb') as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    response_format="text"
                )
            
            result = f"""TIKTOK VIDEO TRANSCRIPT
Title: {title}
Creator: {creator}
URL: {url}

TRANSCRIPT:
{transcript}
"""
            logger.info(f"Successfully transcribed TikTok ({len(transcript)} characters)")
            return result
            
        except Exception as e:
            raise Exception(f"TikTok transcription failed: {e}")


def extract_web_article(url: str) -> str:
    """
    Extract text content from a web article.
    
    Uses BeautifulSoup to parse HTML and extract main content.
    
    Args:
        url: Web page URL
        
    Returns:
        Extracted text content
    """
    try:
        import requests
        from bs4 import BeautifulSoup
    except ImportError:
        raise Exception(
            "Missing dependencies. Install with: "
            "pip install requests beautifulsoup4"
        )
    
    logger.info(f"Extracting content from: {url}")
    
    try:
        # Fetch page
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'
        }
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        
        # Parse HTML
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script in soup(["script", "style", "nav", "footer", "header"]):
            script.decompose()
        
        # Get title
        title = soup.find('title')
        title_text = title.get_text() if title else 'Unknown'
        
        # Extract main content
        # Try common content containers
        main_content = None
        for selector in ['article', 'main', '[role="main"]', '.content', '#content']:
            main_content = soup.select_one(selector)
            if main_content:
                break
        
        if not main_content:
            main_content = soup.find('body')
        
        # Get text
        text = main_content.get_text(separator='\n', strip=True) if main_content else ''
        
        # Clean up whitespace
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        text = '\n'.join(lines)
        
        result = f"""WEB ARTICLE
Title: {title_text}
URL: {url}

CONTENT:
{text}
"""
        
        logger.info(f"Extracted {len(text)} characters")
        return result
        
    except Exception as e:
        raise Exception(f"Failed to extract web content: {e}")


# =============================================================================
# Code File Extensions Map (Phase 1 Expansion)
# =============================================================================

CODE_EXTENSIONS = {
    # Programming languages
    '.py': 'Python',
    '.js': 'JavaScript',
    '.ts': 'TypeScript',
    '.jsx': 'JavaScript React',
    '.tsx': 'TypeScript React',
    '.java': 'Java',
    '.cpp': 'C++',
    '.cc': 'C++',
    '.cxx': 'C++',
    '.c': 'C',
    '.h': 'C/C++ Header',
    '.hpp': 'C++ Header',
    '.go': 'Go',
    '.rs': 'Rust',
    '.rb': 'Ruby',
    '.php': 'PHP',
    '.swift': 'Swift',
    '.kt': 'Kotlin',
    '.scala': 'Scala',
    '.r': 'R',
    '.m': 'MATLAB',

    # Web development
    '.html': 'HTML',
    '.css': 'CSS',
    '.scss': 'SCSS',
    '.sass': 'Sass',
    '.vue': 'Vue',

    # Shell & scripts
    '.sh': 'Shell Script',
    '.bash': 'Bash Script',
    '.zsh': 'Zsh Script',
    '.fish': 'Fish Script',
    '.ps1': 'PowerShell',

    # Data & config
    '.sql': 'SQL',
    '.yaml': 'YAML',
    '.yml': 'YAML',
    '.toml': 'TOML',
    '.xml': 'XML',
    '.ini': 'INI Config',
    '.conf': 'Config File',

    # Documentation
    '.rst': 'reStructuredText',
    '.tex': 'LaTeX',
}


def ingest_upload_file(filename: str, content_bytes: bytes) -> str:
    """
    Process an uploaded file and extract text.

    Supports:
    - PDF files (text extraction)
    - Text files (.txt, .md)
    - Audio files (.mp3, .wav, .m4a) - transcription via Whisper
    - Word documents (.docx)
    - Jupyter notebooks (.ipynb) - Phase 1
    - Code files (Python, JavaScript, etc.) - Phase 1
    - Excel spreadsheets (.xlsx, .xls) - Phase 2
    - PowerPoint presentations (.pptx) - Phase 2
    - ZIP archives (.zip) - Phase 3
    - EPUB books (.epub) - Phase 3
    - Subtitle files (.srt, .vtt) - Phase 3

    ✅ Audio files now compressed if over 25MB limit.

    Args:
        filename: Original filename
        content_bytes: File content as bytes

    Returns:
        Extracted text content
    """
    file_ext = Path(filename).suffix.lower()

    logger.info(f"Processing uploaded file: {filename} ({len(content_bytes)} bytes)")

    # Jupyter notebooks (Phase 1)
    if file_ext == '.ipynb':
        return extract_jupyter_notebook(content_bytes, filename)

    # Code files (Phase 1)
    elif file_ext in CODE_EXTENSIONS:
        return extract_code_file(content_bytes, filename)

    # Text files
    elif file_ext in ['.txt', '.md', '.csv', '.json']:
        try:
            return content_bytes.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return content_bytes.decode('latin-1')
            except:
                raise Exception("Failed to decode text file")

    # PDF files
    elif file_ext == '.pdf':
        return extract_pdf_text(content_bytes)

    # Audio files
    elif file_ext in ['.mp3', '.wav', '.m4a', '.ogg', '.flac']:
        return transcribe_audio_file(content_bytes, filename)

    # Word documents
    elif file_ext == '.docx':
        return extract_docx_text(content_bytes)

    # Excel spreadsheets (Phase 2)
    elif file_ext in ['.xlsx', '.xls']:
        return extract_excel_text(content_bytes, filename)

    # PowerPoint presentations (Phase 2)
    elif file_ext == '.pptx':
        return extract_powerpoint_text(content_bytes, filename)

    # ZIP archives (Phase 3)
    elif file_ext == '.zip':
        return extract_zip_archive(content_bytes, filename)

    # EPUB books (Phase 3)
    elif file_ext == '.epub':
        return extract_epub_text(content_bytes, filename)

    # Subtitle files (Phase 3)
    elif file_ext in ['.srt', '.vtt']:
        return extract_subtitles(content_bytes, filename)

    else:
        raise Exception(f"Unsupported file type: {file_ext}")


def extract_pdf_text(content_bytes: bytes) -> str:
    """Extract text from PDF file."""
    try:
        from pypdf import PdfReader
        import io
    except ImportError:
        raise Exception("Install pypdf: pip install pypdf")
    
    try:
        pdf_file = io.BytesIO(content_bytes)
        reader = PdfReader(pdf_file)
        
        text_parts = []
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            if text.strip():
                text_parts.append(f"--- Page {i+1} ---\n{text}")
        
        result = f"PDF DOCUMENT ({len(reader.pages)} pages)\n\n" + "\n\n".join(text_parts)
        
        logger.info(f"Extracted text from {len(reader.pages)} pages")
        return result
        
    except Exception as e:
        raise Exception(f"PDF extraction failed: {e}")


def transcribe_audio_file(content_bytes: bytes, filename: str) -> str:
    """
    Transcribe an audio file using Whisper.
    
    ✅ Now includes compression for files over 25MB.
    
    Args:
        content_bytes: Audio file content
        filename: Original filename
        
    Returns:
        Transcript with metadata
    """
    try:
        from openai import OpenAI
    except ImportError:
        raise Exception("Install openai: pip install openai")
    
    if not OPENAI_API_KEY:
        raise Exception("OPENAI_API_KEY not set")
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        
        # Save original file
        original_path = temp_path / f"original{Path(filename).suffix}"
        with open(original_path, 'wb') as f:
            f.write(content_bytes)
        
        original_size = len(content_bytes)
        logger.info(f"Audio file size: {original_size/(1024*1024):.2f}MB")
        
        # Compress if needed
        if original_size > WHISPER_MAX_SIZE_BYTES:
            logger.info("Compressing audio file...")
            compressed_path = temp_path / f"compressed{Path(filename).suffix}"
            compress_audio_for_whisper(original_path, compressed_path)
            transcription_path = compressed_path
        else:
            transcription_path = original_path
        
        try:
            client = OpenAI(api_key=OPENAI_API_KEY)
            
            with open(transcription_path, 'rb') as audio_file:
                transcript = client.audio.transcriptions.create(
                    model="whisper-1",
                    file=audio_file,
                    response_format="text"
                )
            
            result = f"""AUDIO FILE TRANSCRIPT
Filename: {filename}

TRANSCRIPT:
{transcript}
"""
            logger.info(f"Successfully transcribed audio file ({len(transcript)} characters)")
            return result
            
        except Exception as e:
            raise Exception(f"Audio transcription failed: {e}")


def extract_docx_text(content_bytes: bytes) -> str:
    """Extract text from Word document."""
    try:
        from docx import Document
        import io
    except ImportError:
        raise Exception("Install python-docx: pip install python-docx")

    try:
        doc_file = io.BytesIO(content_bytes)
        doc = Document(doc_file)

        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        result = "WORD DOCUMENT\n\n" + "\n\n".join(text_parts)

        logger.info(f"Extracted text from Word document: {len(text_parts)} paragraphs")
        return result

    except Exception as e:
        raise Exception(f"Word document extraction failed: {e}")


# =============================================================================
# Phase 1 Expansion: Jupyter Notebooks & Code Files
# =============================================================================

def extract_jupyter_notebook(content_bytes: bytes, filename: str) -> str:
    """
    Extract content from Jupyter notebook (.ipynb).

    Extracts:
    - Code cells with syntax highlighting hints
    - Markdown cells
    - Cell outputs (text)
    - Notebook metadata

    Args:
        content_bytes: Notebook file content
        filename: Original filename

    Returns:
        Formatted text with all notebook content
    """
    import json

    try:
        # Parse notebook JSON
        notebook = json.loads(content_bytes.decode('utf-8'))
    except json.JSONDecodeError as e:
        raise Exception(f"Invalid Jupyter notebook format: {e}")
    except UnicodeDecodeError as e:
        raise Exception(f"Failed to decode notebook: {e}")

    # Extract metadata
    metadata = notebook.get('metadata', {})
    kernel_info = metadata.get('kernelspec', {})
    language = kernel_info.get('language', 'unknown')
    kernel_name = kernel_info.get('display_name', 'Unknown')

    text_parts = [
        f"JUPYTER NOTEBOOK: {filename}",
        f"Kernel: {kernel_name}",
        f"Language: {language}",
        ""
    ]

    # Extract cells
    cells = notebook.get('cells', [])
    code_cell_count = 0
    markdown_cell_count = 0

    for i, cell in enumerate(cells, 1):
        cell_type = cell.get('cell_type', 'unknown')
        source = cell.get('source', [])

        # Handle source as list or string
        if isinstance(source, list):
            source_text = ''.join(source)
        else:
            source_text = source

        if not source_text.strip():
            continue  # Skip empty cells

        if cell_type == 'code':
            code_cell_count += 1
            text_parts.append(f"[Code Cell {code_cell_count}]")
            text_parts.append(source_text)

            # Extract outputs if present
            outputs = cell.get('outputs', [])
            for output in outputs:
                # Text output
                if 'text' in output:
                    output_text = output['text']
                    if isinstance(output_text, list):
                        output_text = ''.join(output_text)
                    text_parts.append(f"[Output]\n{output_text}")

                # Data output (e.g., pandas dataframes)
                elif 'data' in output:
                    data = output['data']
                    if 'text/plain' in data:
                        text_parts.append(f"[Output]\n{data['text/plain']}")

        elif cell_type == 'markdown':
            markdown_cell_count += 1
            text_parts.append(f"[Markdown {markdown_cell_count}]")
            text_parts.append(source_text)

        elif cell_type == 'raw':
            text_parts.append(f"[Raw Cell {i}]")
            text_parts.append(source_text)

    result = "\n\n".join(text_parts)

    logger.info(
        f"Extracted Jupyter notebook: {code_cell_count} code cells, "
        f"{markdown_cell_count} markdown cells, {len(result)} characters"
    )

    return result


def extract_code_file(content_bytes: bytes, filename: str) -> str:
    """
    Extract content from source code file with metadata.

    Provides syntax-aware processing with:
    - Language detection
    - Line count statistics
    - Function/class detection hints
    - Preserved code formatting

    Args:
        content_bytes: Source code content
        filename: Original filename

    Returns:
        Formatted code with metadata
    """
    ext = Path(filename).suffix.lower()
    language = CODE_EXTENSIONS.get(ext, 'Unknown')

    # Decode content
    try:
        code = content_bytes.decode('utf-8')
    except UnicodeDecodeError:
        # Try latin-1 as fallback
        try:
            code = content_bytes.decode('latin-1', errors='replace')
            logger.warning(f"Used latin-1 fallback for {filename}")
        except Exception as e:
            raise Exception(f"Failed to decode {filename}: {e}")

    # Calculate statistics
    lines = code.split('\n')
    total_lines = len(lines)

    # Count non-empty, non-comment lines (rough heuristic)
    if language == 'Python':
        code_lines = [l for l in lines if l.strip() and not l.strip().startswith('#')]
    elif language in ['JavaScript', 'TypeScript', 'Java', 'C', 'C++', 'Go', 'Rust']:
        code_lines = [l for l in lines if l.strip() and not l.strip().startswith('//')]
    else:
        code_lines = [l for l in lines if l.strip()]

    loc = len(code_lines)

    # Detect functions/classes (simple heuristics for common patterns)
    functions = []
    classes = []

    if language == 'Python':
        functions = [l.strip() for l in lines if l.strip().startswith('def ')]
        classes = [l.strip() for l in lines if l.strip().startswith('class ')]
    elif language in ['JavaScript', 'TypeScript']:
        functions = [l.strip() for l in lines if 'function ' in l or '=>' in l]
        classes = [l.strip() for l in lines if l.strip().startswith('class ')]
    elif language in ['Java', 'C++', 'C#']:
        classes = [l.strip() for l in lines if 'class ' in l]

    # Build result
    text_parts = [
        f"SOURCE CODE FILE: {filename}",
        f"Language: {language}",
        f"Total Lines: {total_lines}",
        f"Code Lines: {loc}",
    ]

    if functions:
        text_parts.append(f"Functions/Methods: {len(functions)}")
    if classes:
        text_parts.append(f"Classes: {len(classes)}")

    text_parts.extend([
        "",
        "CODE:",
        code
    ])

    result = "\n".join(text_parts)

    logger.info(
        f"Extracted {language} code: {filename} "
        f"({total_lines} lines, {loc} code lines)"
    )

    return result


# =============================================================================
# Phase 2 Expansion: Office Suite (Excel & PowerPoint)
# =============================================================================

def extract_excel_text(content_bytes: bytes, filename: str) -> str:
    """
    Extract text from Excel spreadsheet (.xlsx, .xls).

    Extracts:
    - All sheets with names
    - Cell values (text, numbers, formulas)
    - Table structure preserved as much as possible
    - Sheet metadata

    Args:
        content_bytes: Excel file content
        filename: Original filename

    Returns:
        Formatted text with all spreadsheet data
    """
    try:
        from openpyxl import load_workbook
        import io
    except ImportError:
        raise Exception("Install openpyxl: pip install openpyxl")

    try:
        # Load workbook from bytes
        wb = load_workbook(io.BytesIO(content_bytes), data_only=True)

        text_parts = [
            f"EXCEL SPREADSHEET: {filename}",
            f"Sheets: {len(wb.sheetnames)}",
            ""
        ]

        total_rows = 0
        total_cells = 0

        # Process each sheet
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]

            # Get sheet dimensions
            max_row = sheet.max_row
            max_col = sheet.max_column

            text_parts.append(f"=== Sheet: {sheet_name} ({max_row} rows × {max_col} cols) ===")
            text_parts.append("")

            # Extract cell values
            sheet_rows = []
            for row in sheet.iter_rows(values_only=True):
                # Convert row to strings, handling None values
                row_values = []
                for cell in row:
                    if cell is None:
                        row_values.append("")
                    elif isinstance(cell, (int, float)):
                        row_values.append(str(cell))
                    else:
                        row_values.append(str(cell))

                # Join with pipe separator for table-like format
                row_text = " | ".join(row_values)

                # Only add non-empty rows
                if row_text.strip(" |"):
                    sheet_rows.append(row_text)
                    total_cells += len([c for c in row_values if c])

            text_parts.extend(sheet_rows)
            text_parts.append("")  # Blank line between sheets

            total_rows += len(sheet_rows)

        # Add summary
        result = "\n".join(text_parts)

        logger.info(
            f"Extracted Excel: {filename} "
            f"({len(wb.sheetnames)} sheets, {total_rows} rows, {total_cells} cells)"
        )

        return result

    except Exception as e:
        raise Exception(f"Excel extraction failed: {e}")


def extract_powerpoint_text(content_bytes: bytes, filename: str) -> str:
    """
    Extract text from PowerPoint presentation (.pptx).

    Extracts:
    - All slides with slide numbers
    - Text from all shapes (titles, content, text boxes)
    - Speaker notes
    - Presentation metadata

    Args:
        content_bytes: PowerPoint file content
        filename: Original filename

    Returns:
        Formatted text with all presentation content
    """
    try:
        from pptx import Presentation
        import io
    except ImportError:
        raise Exception("Install python-pptx: pip install python-pptx")

    try:
        # Load presentation from bytes
        prs = Presentation(io.BytesIO(content_bytes))

        text_parts = [
            f"POWERPOINT PRESENTATION: {filename}",
            f"Slides: {len(prs.slides)}",
            ""
        ]

        total_shapes = 0
        total_notes = 0

        # Process each slide
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {i} ---")
            text_parts.append("")

            slide_texts = []

            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_texts.append(shape.text)
                    total_shapes += 1

                # Handle tables
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        if row_text.strip():
                            slide_texts.append(f"[Table] {row_text}")

            # Add slide content
            if slide_texts:
                text_parts.extend(slide_texts)
            else:
                text_parts.append("[Empty slide]")

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame.text:
                    text_parts.append("")
                    text_parts.append(f"[Speaker Notes]")
                    text_parts.append(notes_frame.text)
                    total_notes += 1

            text_parts.append("")  # Blank line between slides

        result = "\n".join(text_parts)

        logger.info(
            f"Extracted PowerPoint: {filename} "
            f"({len(prs.slides)} slides, {total_shapes} text shapes, {total_notes} notes)"
        )

        return result

    except Exception as e:
        raise Exception(f"PowerPoint extraction failed: {e}")


# ============================================================================
# PHASE 3: ARCHIVES & E-BOOKS
# ============================================================================

def extract_zip_archive(content_bytes: bytes, filename: str) -> str:
    """
    Extract and process ZIP archive contents.

    Recursively processes all supported file types within the archive.
    Skips directories and files larger than 10MB to prevent resource exhaustion.

    Supported files within ZIP:
    - All file types supported by ingest_upload_file()
    - Including code files, Jupyter notebooks, Office docs, PDFs, etc.

    Args:
        content_bytes: Raw bytes of ZIP file
        filename: Original filename for metadata

    Returns:
        Formatted text with metadata and extracted content from all files

    Example Output:
        ZIP ARCHIVE: code_project.zip
        Files: 15 total, 12 processed, 3 skipped
        Total size: 2.4 MB

        === src/main.py ===
        [Python code content]

        === docs/README.md ===
        [Markdown content]
    """
    import zipfile
    import io

    try:
        zip_file = zipfile.ZipFile(io.BytesIO(content_bytes))

        # Gather file statistics
        total_files = 0
        total_size = 0
        processed_files = 0
        skipped_files = 0

        text_parts = []
        text_parts.append(f"ZIP ARCHIVE: {filename}")
        text_parts.append("=" * 60)
        text_parts.append("")

        # First pass: collect statistics
        file_list = []
        for file_info in zip_file.infolist():
            if not file_info.is_dir():
                total_files += 1
                total_size += file_info.file_size
                file_list.append(file_info)

        text_parts.append(f"Total files: {total_files}")
        text_parts.append(f"Total size: {total_size / (1024*1024):.2f} MB")
        text_parts.append("")
        text_parts.append("CONTENTS:")
        text_parts.append("-" * 60)
        text_parts.append("")

        # Second pass: process each file
        for file_info in file_list:
            # Skip large files (> 10MB per file)
            if file_info.file_size > 10 * 1024 * 1024:
                text_parts.append(f"⚠️  SKIPPED (too large): {file_info.filename}")
                text_parts.append(f"   Size: {file_info.file_size / (1024*1024):.2f} MB")
                text_parts.append("")
                skipped_files += 1
                continue

            # Skip hidden files and system files
            if file_info.filename.startswith('.') or '__MACOSX' in file_info.filename:
                skipped_files += 1
                continue

            try:
                file_content = zip_file.read(file_info.filename)

                # Recursively process supported file types
                extracted_text = ingest_upload_file(file_info.filename, file_content)

                text_parts.append(f"=== {file_info.filename} ===")
                text_parts.append(extracted_text)
                text_parts.append("")
                text_parts.append("-" * 60)
                text_parts.append("")

                processed_files += 1

            except Exception as e:
                text_parts.append(f"⚠️  FAILED: {file_info.filename}")
                text_parts.append(f"   Error: {str(e)}")
                text_parts.append("")
                skipped_files += 1

        # Summary at the end
        text_parts.append("")
        text_parts.append("SUMMARY:")
        text_parts.append(f"Processed: {processed_files} files")
        text_parts.append(f"Skipped: {skipped_files} files")
        if total_files > 0:
            text_parts.append(f"Success rate: {(processed_files/total_files*100):.1f}%")
        else:
            text_parts.append(f"Success rate: N/A (empty archive)")

        result = "\n".join(text_parts)

        logger.info(
            f"Extracted ZIP archive: {filename} "
            f"({total_files} files, {processed_files} processed, {skipped_files} skipped)"
        )

        return result

    except zipfile.BadZipFile:
        raise Exception(f"Invalid ZIP file: {filename}")
    except Exception as e:
        raise Exception(f"ZIP extraction failed: {e}")


def extract_epub_text(content_bytes: bytes, filename: str) -> str:
    """
    Extract text from EPUB book.

    EPUB is a popular e-book format used for digital books, technical documentation,
    and educational materials. This function extracts:
    - Book metadata (title, author, language)
    - All chapter content
    - Table of contents structure

    Args:
        content_bytes: Raw bytes of EPUB file
        filename: Original filename for metadata

    Returns:
        Formatted text with metadata and all chapter content

    Example Output:
        EPUB BOOK: Python Programming Guide
        Author: John Doe
        Language: en

        === Chapter 1: Introduction ===
        [Chapter content]

        === Chapter 2: Getting Started ===
        [Chapter content]
    """
    import ebooklib
    from ebooklib import epub
    from bs4 import BeautifulSoup
    import io

    try:
        book = epub.read_epub(io.BytesIO(content_bytes))
        text_parts = []

        # Extract metadata
        title = filename  # Default to filename
        author = "Unknown"
        language = "unknown"

        try:
            if book.get_metadata('DC', 'title'):
                title = book.get_metadata('DC', 'title')[0][0]
        except:
            pass

        try:
            if book.get_metadata('DC', 'creator'):
                author = book.get_metadata('DC', 'creator')[0][0]
        except:
            pass

        try:
            if book.get_metadata('DC', 'language'):
                language = book.get_metadata('DC', 'language')[0][0]
        except:
            pass

        text_parts.append(f"EPUB BOOK: {title}")
        text_parts.append("=" * 60)
        text_parts.append(f"Author: {author}")
        text_parts.append(f"Language: {language}")
        text_parts.append(f"Filename: {filename}")
        text_parts.append("")
        text_parts.append("CONTENT:")
        text_parts.append("-" * 60)
        text_parts.append("")

        # Extract chapters
        chapter_count = 0
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                try:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')

                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()

                    text = soup.get_text(separator='\n', strip=True)

                    if text and len(text.strip()) > 10:  # Skip very short sections (reduced threshold)
                        chapter_count += 1

                        # Try to extract chapter title
                        chapter_title = f"Chapter {chapter_count}"
                        h1 = soup.find('h1')
                        if h1:
                            chapter_title = h1.get_text(strip=True)

                        text_parts.append(f"=== {chapter_title} ===")
                        text_parts.append("")
                        text_parts.append(text)
                        text_parts.append("")
                        text_parts.append("-" * 60)
                        text_parts.append("")

                except Exception as e:
                    logger.warning(f"Failed to extract EPUB chapter: {e}")
                    continue

        text_parts.append("")
        text_parts.append(f"Total chapters extracted: {chapter_count}")

        result = "\n".join(text_parts)

        logger.info(
            f"Extracted EPUB: {title} by {author} "
            f"({chapter_count} chapters)"
        )

        return result

    except Exception as e:
        raise Exception(f"EPUB extraction failed: {e}")


def extract_subtitles(content_bytes: bytes, filename: str) -> str:
    """
    Extract text from subtitle files (SRT, VTT).

    Subtitle files are commonly used for:
    - Video transcripts
    - Language learning materials
    - Accessibility content
    - Lecture transcriptions

    Supported formats:
    - SRT (SubRip): Most common format, includes timestamps
    - VTT (WebVTT): Web-based subtitle format

    Args:
        content_bytes: Raw bytes of subtitle file
        filename: Original filename for metadata

    Returns:
        Formatted text with all subtitle content (timestamps removed)

    Example Output:
        SUBTITLE FILE: lecture_01.srt
        Format: SRT
        Entries: 145

        Hello and welcome to today's lecture.
        In this session we'll cover...
    """
    from pathlib import Path

    ext = Path(filename).suffix.lower()
    text_parts = []

    try:
        if ext == '.srt':
            # SRT format (manual parsing - simple format, no library needed)
            text = content_bytes.decode('utf-8')
            lines = text.split('\n')

            text_parts.append(f"SUBTITLE FILE: {filename}")
            text_parts.append("=" * 60)
            text_parts.append(f"Format: SRT (SubRip)")

            # Parse SRT manually
            subtitle_texts = []
            current_subtitle = []
            in_subtitle_text = False

            for line in lines:
                line_stripped = line.strip()

                # Skip empty lines
                if not line_stripped:
                    if current_subtitle:
                        subtitle_texts.append('\n'.join(current_subtitle))
                        current_subtitle = []
                    in_subtitle_text = False
                    continue

                # Skip subtitle numbers (lines that are just digits)
                if line_stripped.isdigit():
                    continue

                # Skip timestamp lines (contain -->)
                if '-->' in line_stripped:
                    in_subtitle_text = True
                    continue

                # This is subtitle text
                if in_subtitle_text or current_subtitle:
                    current_subtitle.append(line_stripped)

            # Don't forget last subtitle
            if current_subtitle:
                subtitle_texts.append('\n'.join(current_subtitle))

            text_parts.append(f"Entries: {len(subtitle_texts)}")
            text_parts.append("")
            text_parts.append("TRANSCRIPT:")
            text_parts.append("-" * 60)
            text_parts.append("")

            # Add all subtitle text
            text_parts.extend(subtitle_texts)

            logger.info(f"Extracted SRT subtitles: {filename} ({len(subtitle_texts)} entries)")

        elif ext == '.vtt':
            # WebVTT format (simple parsing)
            text = content_bytes.decode('utf-8')
            lines = text.split('\n')

            text_parts.append(f"SUBTITLE FILE: {filename}")
            text_parts.append("=" * 60)
            text_parts.append(f"Format: WebVTT")
            text_parts.append("")
            text_parts.append("TRANSCRIPT:")
            text_parts.append("-" * 60)
            text_parts.append("")

            # Filter out WebVTT header, timestamps, and metadata
            subtitle_lines = []
            for line in lines:
                line = line.strip()
                # Skip empty lines, WebVTT header, timestamps, and cue settings
                if (line and
                    not line.startswith('WEBVTT') and
                    '-->' not in line and
                    not line.startswith('NOTE') and
                    not line.isdigit()):  # Skip cue numbers
                    subtitle_lines.append(line)

            text_parts.extend(subtitle_lines)

            logger.info(f"Extracted VTT subtitles: {filename} ({len(subtitle_lines)} lines)")

        else:
            raise Exception(f"Unsupported subtitle format: {ext}")

        result = "\n".join(text_parts)
        return result

    except Exception as e:
        raise Exception(f"Subtitle extraction failed: {e}")
